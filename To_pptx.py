# beamerなどの横向きのpdfファイルをパワーポイントファイルに変換.


import os
import sys
import shutil
from pathlib import Path
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches


# （参考）https://zenn.dev/torachi0401/articles/change_pdf_to_image
def pdf_to_image(path):
  # poppler/binを環境変数PATHに追加する
  poppler_dir = Path(__file__).parent.absolute() / "poppler/Library/bin"
  os.environ["PATH"] += os.pathsep + str(poppler_dir)

  pdf_path = Path(path)
  image_dir = Path(pdf_path.parent, "images") 

  if not os.path.isdir(image_dir):
    os.mkdir(image_dir)
  else:
    # 画像ファイル格納フォルダがすでに存在する場合、エラーとする
    raise Exception(f"Directory already exists: {image_dir}")

  # PDFファイルを画像に変換
  pages = convert_from_path(pdf_path, dpi=400)
  if len(pages) > 1:
      for i, page in enumerate(pages):
          file_name = pdf_path.stem + "_{:02d}".format(i + 1) + ".jpeg"
          image_path = image_dir / file_name
          # JPEGで保存
          page.save(str(image_path), "JPEG")
  else:
      file_name = pdf_path.stem + ".jpeg"
      image_path = image_dir / file_name
      # JPEGで保存
      pages[0].save(str(image_path), "JPEG")
  
  return str(image_dir)





def images_to_pptx(image_dir, output_pptx):
  # PowerPointプレゼンテーションを作成
  prs = Presentation()

  # 画像フォルダ内のファイルを取得（ソートして順番を維持）
  images = sorted([f for f in os.listdir(image_dir) if f.lower().endswith(('.png', '.jpg', '.jpeg', '.bmp', '.gif'))])

  for image in images:
    img_path = os.path.join(image_dir, image)
    
    # スライドを追加
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # 白紙スライド
    
    # 画像をスライドの中央に配置
    left = top = Inches(0)  # 余白
    slide_width = prs.slide_width - 2 * left
    slide_height = prs.slide_height - 2 * top

    slide.shapes.add_picture(img_path, left, top, slide_width, slide_height)

  # PPTXファイルを保存
  prs.save(output_pptx)






def main():
  try:
    # PDFファイルパスをコマンドライン引数で受け取る
    pdf_path = sys.argv[1]

    pptx_name = os.path.splitext(os.path.basename(pdf_path))[0] + ".pptx"
    pptx_path = os.path.join(os.path.dirname(pdf_path), pptx_name)

    # PDFを画像ファイルに変換
    image_dir = pdf_to_image(pdf_path)

    # 画像ファイルをpptxに変換
    images_to_pptx(image_dir, pptx_path)

    print(f"PowerPoint file has been created: {pptx_path}")
    
    #不要になった画像フォルダを消去
    shutil.rmtree(image_dir)

  except Exception as e:
    print(e)
    print("Failed to convert to PowerPoint file.")




if __name__ == '__main__':
  main()
  
  
  